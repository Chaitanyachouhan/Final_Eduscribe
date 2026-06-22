from pptx import Presentation

def extract_ppt_text(ppt_path: str) -> list[str]:
    prs = Presentation(ppt_path)
    slide_texts = []
    for slide in prs.slides:
        text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slide_texts.append(text.strip())
    return slide_texts